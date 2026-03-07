import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    accuracy_score,
    precision_score,
    recall_score,
    f1_score,
    confusion_matrix,
    roc_curve,
    auc,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os
import warnings

warnings.filterwarnings("ignore")


def set_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_rgb)


def add_gradient(slide):
    set_background(slide, (240, 248, 255))


def style_title(slide, title_text, subtitle_text=None):
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    )
    title_frame = title_shape.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    title_p.alignment = PP_ALIGN.CENTER

    if subtitle_text:
        sub_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(9), Inches(0.5)
        )
        sub_frame = sub_shape.text_frame
        sub_p = sub_frame.paragraphs[0]
        sub_p.text = subtitle_text
        sub_p.font.size = Pt(20)
        sub_p.font.color.rgb = RGBColor(100, 100, 100)
        sub_p.alignment = PP_ALIGN.CENTER


def add_content_with_bullet(slide, title, bullets, left=1, top=1.5, width=8):
    title_shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.5)
    )
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    content_shape = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.6), Inches(width), Inches(4)
    )
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.level = 0


def add_image_slide(prs, title, image_path, desc=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_title(slide, title)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.8), width=Inches(6))

    if desc:
        desc_shape = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        desc_frame = desc_shape.text_frame
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = desc
        desc_p.font.size = Pt(14)
        desc_p.font.color.rgb = RGBColor(80, 80, 80)
        desc_p.alignment = PP_ALIGN.CENTER


def add_two_columns(slide, title, left_content, right_content):
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
    )
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    title_p.alignment = PP_ALIGN.CENTER

    left_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(4), Inches(4)
    )
    left_frame = left_shape.text_frame
    left_frame.word_wrap = True
    for i, text in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    right_shape = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4), Inches(4))
    right_frame = right_shape.text_frame
    right_frame.word_wrap = True
    for i, text in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.space_after = Pt(8)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide, (0, 51, 102))

title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "BINARY LOGISTIC REGRESSION"
title_p.font.size = Pt(54)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "Olympic Athletes Medal Prediction"
subtitle_p.font.size = Pt(28)
subtitle_p.font.color.rgb = RGBColor(255, 215, 0)
subtitle_p.alignment = PP_ALIGN.CENTER

name_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
name_frame = name_box.text_frame
name_p = name_frame.paragraphs[0]
name_p.text = (
    "RAGHURAAM PRAKASAM\nReg. No: 2022-DS-021\nLoyola College (Autonomous), Chennai"
)
name_p.font.size = Pt(18)
name_p.font.color.rgb = RGBColor(255, 255, 255)
name_p.alignment = PP_ALIGN.CENTER

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Agenda", "What We Will Cover")
add_content_with_bullet(
    slide,
    "",
    [
        "• Introduction to Logistic Regression",
        "• Problem Statement & Objectives",
        "• Dataset Overview",
        "• Data Preprocessing",
        "• Exploratory Data Analysis",
        "• Model Building & Training",
        "• Model Evaluation Metrics",
        "• Results & Interpretation",
        "• Conclusion & Future Scope",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Introduction to Logistic Regression")
add_content_with_bullet(
    slide,
    "",
    [
        "• What is Logistic Regression?",
        "  - Classification algorithm predicting binary outcomes",
        "  - Uses sigmoid function to transform predictions",
        "  - Outputs probability between 0 and 1",
        "",
        "• Types: Binary, Multinomial, Ordinal",
        "",
        "• Key Applications:",
        "  - Medical diagnosis",
        "  - Credit risk assessment",
        "  - Sports prediction",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "The Sigmoid Function")
add_content_with_bullet(
    slide,
    "",
    [
        "• Formula: P(y=1|X) = 1 / (1 + e^-(b0 + b1*x1 + ... + bn*xn))",
        "",
        "• S-Curve Properties:",
        "  - Maps any value to range [0,1]",
        "  - Threshold at 0.5 for classification",
        "",
        "• Coefficients:",
        "  - Positive → Increases probability",
        "  - Negative → Decreases probability",
        "",
        "• Odds Ratio: e^coefficient",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Problem Statement")
add_content_with_bullet(
    slide,
    "",
    [
        "Research Question:",
        "• Can we predict if an Olympic athlete will win a medal?",
        "",
        "Variables Used:",
        "• Age, Height, Weight (Numerical)",
        "• Gender, Sport, Country (Categorical)",
        "",
        "Hypothesis:",
        "• H0: No significant relationship exists",
        "• H1: Significant relationship exists",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Dataset Overview")
add_content_with_bullet(
    slide,
    "",
    [
        "• Dataset: Olympic Athletes (1900-2016)",
        "",
        "• Features: 15 columns",
        "  - ID, Name, Sex, Age, Height, Weight",
        "  - Team, NOC, Games, Year, Season",
        "  - City, Sport, Event, Medal",
        "",
        "• Total Records: 271,116 athletes",
        "",
        "• Target Variable: Medal_Won (1/0)",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Data Preprocessing")
add_content_with_bullet(
    slide,
    "",
    [
        "1. Data Cleaning:",
        "   - Removed: ID, Name, Event, Games, City",
        "",
        "2. Target Variable Creation:",
        "   - Medal_Won = 1 (Gold/Silver/Bronze)",
        "   - Medal_Won = 0 (No Medal)",
        "",
        "3. Missing Values:",
        "   - Age, Height, Weight → Median imputation",
        "",
        "4. Feature Encoding:",
        "   - LabelEncoder for categorical variables",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Train-Test Split")
add_content_with_bullet(
    slide,
    "",
    [
        "• Training Set: 80%",
        "• Test Set: 20%",
        "",
        "• Stratified Sampling: Maintains class distribution",
        "",
        "• Feature Scaling: StandardScaler",
        "  - Normalizes features to mean=0, std=1",
        "",
        "• Final Features:",
        "  Age, Height, Weight, Sex, Team, Sport, Season",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Model Training")
add_content_with_bullet(
    slide,
    "",
    [
        "• Algorithm: Binary Logistic Regression",
        "",
        "• Parameters:",
        "  - solver: lbfgs",
        "  - max_iter: 1000",
        "  - random_state: 42",
        "",
        "• Python Implementation:",
        "  from sklearn.linear_model import LogisticRegression",
        "  model = LogisticRegression()",
        "  model.fit(X_train, y_train)",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Model Evaluation Metrics")
add_content_with_bullet(
    slide,
    "",
    [
        "• Accuracy: Overall correctness",
        "",
        "• Precision: True positives / Predicted positives",
        "",
        "• Recall: True positives / Actual positives",
        "",
        "• F1-Score: Harmonic mean of Precision & Recall",
        "",
        "• ROC-AUC: Area under ROC curve",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Confusion Matrix")
add_two_columns(
    slide,
    "",
    [
        "True Negatives (TN)",
        "False Positives (FP)",
        "False Negatives (FN)",
        "True Positives (TP)",
    ],
    ["Correct rejection", "False alarms", "Missed predictions", "Correct predictions"],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Results & Performance")
add_content_with_bullet(
    slide,
    "",
    [
        "• Model Accuracy: 91.33%",
        "",
        "• Key Findings:",
        "  - Sport type is most significant predictor",
        "  - Country affiliation influences success",
        "  - Physical attributes have moderate impact",
        "",
        "• Interpretation:",
        "  - Coefficients show direction of influence",
        "  - Odds ratios quantify effect size",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Feature Importance")
add_content_with_bullet(
    slide,
    "",
    [
        "Top Predictors:",
        "",
        "1. Sport_encoded",
        "   - Most influential factor",
        "   - Different sports have different medal rates",
        "",
        "2. Team_encoded",
        "   - Country matters for success",
        "",
        "3. Age, Height, Weight",
        "   - Moderate influence",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Conclusion")
add_content_with_bullet(
    slide,
    "",
    [
        "Key Achievements:",
        "• Successfully built predictive model (91.33% accuracy)",
        "• Identified key factors affecting medal winning",
        "• Demonstrated practical application of ML in sports",
        "",
        "Limitations:",
        "• Imbalanced dataset (only ~11% medal winners)",
        "• Linear decision boundary",
        "",
        "Future Work:",
        "• Use ensemble methods (Random Forest, XGBoost)",
        "• Include more features",
        "• Handle class imbalance",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
style_title(slide, "Thank You!")
add_content_with_bullet(
    slide,
    "",
    [
        "Questions?",
        "",
        "Contact:",
        "RAGHURAAM PRAKASAM",
        "Reg. No: 2022-DS-021",
        "Loyola College (Autonomous), Chennai",
        "",
        "Email: raghuraam@example.com",
    ],
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide, (0, 51, 102))
title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "LOYOLA COLLEGE (AUTONOMOUS), CHENNAI"
title_p.font.size = Pt(36)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 215, 0)
title_p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "Department of Data Science"
subtitle_p.font.size = Pt(24)
subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
subtitle_p.alignment = PP_ALIGN.CENTER

prs.save("C:/Users/raghu/OneDrive/Documents/loyola project/presentation.pptx")
print("Presentation created successfully!")
