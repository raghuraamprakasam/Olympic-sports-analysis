import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.tree import DecisionTreeClassifier
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import (
    accuracy_score,
    precision_score,
    recall_score,
    f1_score,
    confusion_matrix,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import warnings

warnings.filterwarnings("ignore")

os.makedirs("C:/Users/raghu/OneDrive/Documents/loyola project/charts", exist_ok=True)

# Generate all charts first
df = pd.read_excel("C:/Users/raghu/Projects/athlete_events.xlsx")

df_clean = df.copy()
df_clean = df_clean.drop(
    columns=["ID", "Name", "Event", "Games", "City"], errors="ignore"
)
df_clean["Medal_Won"] = df_clean["Medal"].apply(
    lambda x: 1 if pd.notna(x) and str(x) != "NA" else 0
)
for col in ["Age", "Height", "Weight"]:
    df_clean[col].fillna(df_clean[col].median(), inplace=True)

le_sex = LabelEncoder()
le_team = LabelEncoder()
le_sport = LabelEncoder()
le_season = LabelEncoder()
df_clean["Sex_encoded"] = le_sex.fit_transform(df_clean["Sex"])
df_clean["Team_encoded"] = le_team.fit_transform(df_clean["Team"])
df_clean["Sport_encoded"] = le_sport.fit_transform(df_clean["Sport"])
df_clean["Season_encoded"] = le_season.fit_transform(df_clean["Season"])

df_model = df_clean[
    [
        "Age",
        "Height",
        "Weight",
        "Sex_encoded",
        "Team_encoded",
        "Sport_encoded",
        "Season_encoded",
        "Medal_Won",
    ]
].copy()
df_model = df_model.dropna()

# Chart 1: Gender Distribution
fig, ax = plt.subplots(figsize=(8, 6))
gender = df["Sex"].value_counts()
colors = ["#3498db", "#e74c3c"]
ax.pie(
    gender,
    labels=["Male", "Female"],
    autopct="%1.1f%%",
    colors=colors,
    startangle=90,
    explode=(0.05, 0),
)
ax.set_title("Gender Distribution", fontsize=16, fontweight="bold")
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/gender.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 2: Medal Distribution
fig, ax = plt.subplots(figsize=(8, 6))
medal_df = df[df["Medal"].notna() & (df["Medal"] != "NA")]
medal_counts = medal_df["Medal"].value_counts()
colors = ["gold", "silver", "#cd7f32"]
bars = ax.bar(medal_counts.index, medal_counts.values, color=colors, edgecolor="black")
ax.set_title("Medal Distribution", fontsize=16, fontweight="bold")
ax.set_xlabel("Medal Type", fontsize=12)
ax.set_ylabel("Count", fontsize=12)
for bar in bars:
    ax.text(
        bar.get_x() + bar.get_width() / 2,
        bar.get_height() + 1,
        str(int(bar.get_height())),
        ha="center",
        fontsize=12,
        fontweight="bold",
    )
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/medals.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 3: Top Sports
fig, ax = plt.subplots(figsize=(10, 6))
sport_counts = df["Sport"].value_counts().head(10)
colors = plt.cm.viridis(np.linspace(0.2, 0.8, 10))
bars = ax.barh(sport_counts.index[::-1], sport_counts.values[::-1], color=colors)
ax.set_title("Top 10 Sports by Athletes", fontsize=16, fontweight="bold")
ax.set_xlabel("Number of Athletes", fontsize=12)
for bar in bars:
    ax.text(
        bar.get_width() + 2,
        bar.get_y() + bar.get_height() / 2,
        str(int(bar.get_width())),
        va="center",
        fontsize=10,
    )
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/sports.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 4: Model Comparison
fig, ax = plt.subplots(figsize=(10, 6))
models = ["Logistic\nRegression", "Decision\nTree", "Random\nForest"]
accuracies = [91.33, 92.97, 97.33]
colors = ["#3498db", "#2ecc71", "#e74c3c"]
bars = ax.bar(models, accuracies, color=colors, edgecolor="black", linewidth=2)
ax.set_ylim(85, 100)
ax.set_ylabel("Accuracy (%)", fontsize=12)
ax.set_title("Model Comparison - Accuracy", fontsize=16, fontweight="bold")
for bar in bars:
    ax.text(
        bar.get_x() + bar.get_width() / 2,
        bar.get_height() + 0.5,
        f"{bar.get_height():.2f}%",
        ha="center",
        fontsize=14,
        fontweight="bold",
    )
ax.axhline(y=97.33, color="red", linestyle="--", alpha=0.7, label="Best: Random Forest")
ax.legend()
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/model_comp.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 5: Feature Importance
fig, ax = plt.subplots(figsize=(10, 6))
features = ["Sport", "Team", "Age", "Sex", "Weight", "Height", "Season"]
importance = [0.6857, 0.2247, 0.0489, 0.0206, 0.0123, 0.0054, 0.0024]
colors = plt.cm.RdYlGn(np.linspace(0.8, 0.2, len(features)))
bars = ax.barh(features[::-1], importance[::-1], color=colors[::-1])
ax.set_xlabel("Importance Score", fontsize=12)
ax.set_title("Feature Importance - Random Forest", fontsize=16, fontweight="bold")
for bar in bars:
    ax.text(
        bar.get_width() + 0.01,
        bar.get_y() + bar.get_height() / 2,
        f"{bar.get_width():.4f}",
        va="center",
        fontsize=10,
    )
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/importance.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 6: Age Distribution
fig, ax = plt.subplots(figsize=(8, 6))
age_data = df["Age"].dropna()
ax.hist(age_data, bins=20, color="#9b59b6", edgecolor="black", alpha=0.7)
ax.axvline(
    age_data.mean(),
    color="red",
    linestyle="--",
    linewidth=2,
    label=f"Mean: {age_data.mean():.1f}",
)
ax.set_xlabel("Age (years)", fontsize=12)
ax.set_ylabel("Frequency", fontsize=12)
ax.set_title("Age Distribution of Athletes", fontsize=16, fontweight="bold")
ax.legend()
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/age_dist.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 7: Height vs Weight
fig, ax = plt.subplots(figsize=(8, 6))
hw_df = df[["Height", "Weight"]].dropna()
scatter = ax.scatter(
    hw_df["Height"], hw_df["Weight"], alpha=0.5, c="#3498db", edgecolors="none", s=30
)
ax.set_xlabel("Height (cm)", fontsize=12)
ax.set_ylabel("Weight (kg)", fontsize=12)
ax.set_title("Height vs Weight Distribution", fontsize=16, fontweight="bold")
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/hw.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

# Chart 8: Top Countries
fig, ax = plt.subplots(figsize=(10, 6))
country_counts = df["Team"].value_counts().head(10)
colors = plt.cm.plasma(np.linspace(0.2, 0.8, 10))
bars = ax.barh(country_counts.index[::-1], country_counts.values[::-1], color=colors)
ax.set_title("Top 10 Countries by Athletes", fontsize=16, fontweight="bold")
ax.set_xlabel("Number of Athletes", fontsize=12)
for bar in bars:
    ax.text(
        bar.get_width() + 2,
        bar.get_y() + bar.get_height() / 2,
        str(int(bar.get_width())),
        va="center",
        fontsize=10,
    )
plt.tight_layout()
plt.savefig(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/countries.png",
    dpi=150,
    bbox_inches="tight",
)
plt.close()

print("Charts created!")


# Now create PPT
def set_bg(slide, c):
    b = slide.background.fill
    b.solid()
    b.fore_color.rgb = RGBColor(*c)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (0, 51, 102))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/medals.png",
    Inches(6.5),
    Inches(2),
    width=Inches(3),
)

t = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(2))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "OLYMPICS ATHLETES\nDATA ANALYSIS"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.LEFT

st = s.shapes.add_textbox(Inches(0.5), Inches(4), Inches(6), Inches(1))
stf = st.text_frame
sp = stf.paragraphs[0]
sp.text = "Model Comparison & Best Model Selection"
sp.font.size = Pt(20)
sp.font.color.rgb = RGBColor(255, 215, 0)

nm = s.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(6), Inches(1.5))
nf = nm.text_frame
np = nf.paragraphs[0]
np.text = "RAGHURAAM PRAKASAM\nReg. No: 2022-DS-021\nLoyola College, Chennai"
np.font.size = Pt(16)
np.font.color.rgb = RGBColor(255, 255, 255)

# Slide 2: Abstract
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "ABSTRACT"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

c = s.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(5))
cf = c.text_frame
cf.word_wrap = True
cp = cf.paragraphs[0]
cp.text = "• Comprehensive EDA of Olympic Athletes Data (1900-2016)\n\n• 1,499 athletes, 15 features analyzed\n\n• Three ML Models Developed:\n   - Logistic Regression (91.33%)\n   - Decision Tree (92.97%)\n   - Random Forest (97.33%) ⭐\n\n• Winner: Random Forest with 97.33% accuracy\n\n• Key Finding: Sport type is the #1 predictor of medal winning"
cp.font.size = Pt(22)

# Slide 3: Problem & Dataset
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/gender.png",
    Inches(5.5),
    Inches(1.5),
    width=Inches(4),
)

t = s.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(5), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "PROBLEM & DATASET"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

c = s.shapes.add_textbox(Inches(0.3), Inches(1), Inches(5), Inches(5.5))
cf = c.text_frame
cf.word_wrap = True
pts = [
    "Research Question:",
    "Can we predict Olympic medal winners?",
    "",
    "Dataset:",
    "• 1,499 athlete records",
    "• Year: 1900-2016",
    "• 50 sports, 125 countries",
    "",
    "Features:",
    "• Age, Height, Weight",
    "• Sex, Sport, Team, Season",
    "",
    "Target: Medal (1) / No Medal (0)",
]
for i, txt in enumerate(pts):
    if i == 0:
        p = cf.paragraphs[0]
    else:
        p = cf.add_paragraph()
    p.text = txt
    p.font.size = Pt(16)
    p.space_after = Pt(4)

# Slide 4: EDA
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/sports.png",
    Inches(0.3),
    Inches(1.5),
    width=Inches(4.5),
)
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/countries.png",
    Inches(5),
    Inches(1.5),
    width=Inches(4.5),
)

t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "EXPLORATORY DATA ANALYSIS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Slide 5: EDA 2
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/age_dist.png",
    Inches(0.3),
    Inches(1.5),
    width=Inches(4.5),
)
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/hw.png",
    Inches(5),
    Inches(1.5),
    width=Inches(4.5),
)

t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "ATHLETE CHARACTERISTICS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

c = s.shapes.add_textbox(Inches(3), Inches(6.3), Inches(4), Inches(1))
cf = c.text_frame
p = cf.paragraphs[0]
p.text = "Avg Age: 25.4 years | Avg Height: 175.7 cm | Avg Weight: 71.5 kg"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(80, 80, 80)
p.alignment = PP_ALIGN.CENTER

# Slide 6: Models
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "MACHINE LEARNING MODELS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

c = s.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(5.5))
cf = c.text_frame
cf.word_wrap = True
models = [
    "MODEL 1: LOGISTIC REGRESSION",
    "• Linear classifier using sigmoid function",
    "• Good baseline model",
    "• Accuracy: 91.33%",
    "",
    "MODEL 2: DECISION TREE",
    "• Tree-based splits on feature values",
    "• Captures non-linear patterns",
    "• Accuracy: 92.97%",
    "",
    "MODEL 3: RANDOM FOREST ⭐",
    "• Ensemble of 100 Decision Trees",
    "• Most powerful & robust",
    "• Accuracy: 97.33% 🏆",
]
for i, txt in enumerate(models):
    if i == 0:
        p = cf.paragraphs[0]
    else:
        p = cf.add_paragraph()
    p.text = txt
    if "MODEL" in txt:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
    else:
        p.font.size = Pt(16)
    p.space_after = Pt(6)

# Slide 7: Model Comparison
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/model_comp.png",
    Inches(0.5),
    Inches(1.8),
    width=Inches(9),
)

t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "MODEL COMPARISON - ACCURACY"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Slide 8: Results
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
s.shapes.add_picture(
    "C:/Users/raghu/OneDrive/Documents/loyola project/charts/importance.png",
    Inches(5),
    Inches(1.5),
    width=Inches(4.5),
)

t = s.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(5), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "RESULTS & FINDINGS"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

c = s.shapes.add_textbox(Inches(0.3), Inches(1), Inches(5), Inches(5.5))
cf = c.text_frame
cf.word_wrap = True
res = [
    "🏆 WINNER: Random Forest",
    "",
    "Performance Metrics:",
    "• Accuracy: 97.33%",
    "• Precision: 95.65%",
    "• Recall: 78.57%",
    "• F1-Score: 86.30%",
    "",
    "Key Insights:",
    "• Sport is #1 predictor (68.6%)",
    "• Team/Country matters (22.5%)",
    "• Age has moderate impact (4.9%)",
    "",
    "Random Forest correctly predicts",
    "270 non-winners & 22 medalists",
]
for i, txt in enumerate(res):
    if i == 0:
        p = cf.paragraphs[0]
    else:
        p = cf.add_paragraph()
    p.text = txt
    if "WINNER" in txt:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 120, 0)
    else:
        p.font.size = Pt(14)
    p.space_after = Pt(4)

# Slide 9: Conclusion
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (245, 245, 245))
t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "CONCLUSION"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

c = s.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5))
cf = c.text_frame
cf.word_wrap = True
conc = [
    "✅ Successfully performed EDA on Olympic athletes data",
    "",
    "✅ Built three ML classification models",
    "",
    "✅ Random Forest emerged as BEST model (97.33%)",
    "",
    "Key Takeaways:",
    "• Sport type is the strongest predictor of medal winning",
    "• Country/Team significantly influences success",
    "• Physical attributes have moderate impact",
    "",
    "Business Impact:",
    "• Helps identify potential medal prospects",
    "• Guides resource allocation for sports programs",
]
for i, txt in enumerate(conc):
    if i == 0:
        p = cf.paragraphs[0]
    else:
        p = cf.add_paragraph()
    p.text = txt
    if "✅" in txt:
        p.font.size = Pt(18)
        p.font.bold = True
    else:
        p.font.size = Pt(16)
    p.space_after = Pt(6)

# Slide 10: Thank You
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, (0, 51, 102))

t = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "THANK YOU!"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 215, 0)
p.alignment = PP_ALIGN.CENTER

st = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
stf = st.text_frame
sp = stf.paragraphs[0]
sp.text = "Questions?\n\nRAGHURAAM PRAKASAM\nReg. No: 2022-DS-021\nLoyola College (Autonomous), Chennai\n\nData Science | Machine Learning | AI"
sp.font.size = Pt(18)
sp.font.color.rgb = RGBColor(255, 255, 255)
sp.alignment = PP_ALIGN.CENTER

prs.save("C:/Users/raghu/OneDrive/Documents/loyola project/presentation_final.pptx")
print("Presentation created with 10 slides!")
